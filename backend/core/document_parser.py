import os
import re
from pathlib import Path


CHUNK_SIZE = 500      # CJK 目标字符数（中文字≈token）
CHUNK_OVERLAP = 50

# 句末标点（中英文）：在其后切句，标点保留在前句末尾
_SENT_END_RE = re.compile(r'(?<=[。！？；!?;])')

# 仅匹配“整个正文都是占位符”的章节。这里刻意保持集合很小，避免把“没有权限”
# 之类包含“无/没有”的有效短句误判为无内容。
_PLACEHOLDER_SECTION_VALUES = frozenset({
    "无",
    "暂无",
    "无内容",
    "不涉及",
    "n/a",
    "-",
})


def _is_placeholder_only_section(text: str) -> bool:
    """判断 Markdown 正文是否完全由低信息占位符组成。

    支持模板里常见的 ``> 无`` 引用写法，但主动保护围栏代码、缩进代码和
    Markdown 表格。函数只做精确整行匹配，不按包含关系或文本长度推断，避免
    误删“测试”“同意”“OK”等短但有效的知识。
    """
    raw_lines = [line for line in text.splitlines() if line.strip()]
    if not raw_lines:
        return False

    # 代码和表格即使内容看似占位符，也可能是需要原样说明的配置/示例。
    if any(
        re.match(r'^\s*(```|~~~)', line)
        or line.startswith("    ")
        or line.startswith("\t")
        or re.match(r'^\s*\|.*\|\s*$', line)
        for line in raw_lines
    ):
        return False

    normalized: list[str] = []
    for line in raw_lines:
        # Markdown 模板常用 blockquote 包裹占位内容，例如 ``> 无``。
        value = re.sub(r'^(?:\s*>\s*)+', '', line).strip().casefold()
        if not value:
            continue
        normalized.append(value)

    return bool(normalized) and all(
        value in _PLACEHOLDER_SECTION_VALUES for value in normalized
    )


def _is_cjk(text: str) -> bool:
    """按中日韩字符占比判定（>20% 视为 CJK），比"字符数/词数"更准，避免英文被误判。"""
    stripped = [ch for ch in text if not ch.isspace()]
    if not stripped:
        return True
    cjk = sum(
        1 for ch in stripped
        if '一' <= ch <= '鿿'    # 中日韩统一表意文字
        or '぀' <= ch <= 'ヿ'    # 日文假名
        or '가' <= ch <= '힣'    # 韩文谚文
    )
    return cjk / len(stripped) > 0.2


def _hard_split(s: str, size: int) -> list[str]:
    """超长且无标点可切的文本兜底按定长硬切（块间重叠由分组层统一处理）。"""
    return [s[i:i + size] for i in range(0, len(s), size) if s[i:i + size].strip()]


def _to_sentences(text: str) -> list[str]:
    """按行 + 中英文句末标点切成句子单元（去掉空行/空句）。"""
    units: list[str] = []
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue
        for part in _SENT_END_RE.split(line):
            part = part.strip()
            if part:
                units.append(part)
    return units


def _split_text(text: str) -> list[str]:
    """句子感知的递归分块：优先在句末/换行处断开，按目标大小打包并带句子级重叠，
    避免从句子中间切断；只有超长且无标点的句子才退化为定长硬切。"""
    text = re.sub(r'\n{3,}', '\n\n', text).strip()
    if not text:
        return []

    cjk = _is_cjk(text)
    sep = "" if cjk else " "
    # 拉丁文每 token 约 4 字符，按字符切时放大目标，与中文(字≈token)大致对齐
    size = CHUNK_SIZE if cjk else CHUNK_SIZE * 4
    overlap = CHUNK_OVERLAP if cjk else CHUNK_OVERLAP * 4

    # 1) 切句；过长的句子先硬切成 <= size 的单元
    units: list[str] = []
    for s in _to_sentences(text):
        units.extend(_hard_split(s, size) if len(s) > size else [s])
    if not units:
        return []

    # 2) 贪心打包成不超过 size 的分组
    groups: list[list[str]] = []
    cur: list[str] = []
    cur_len = 0
    for u in units:
        if cur and cur_len + len(u) > size:
            groups.append(cur)
            cur, cur_len = [], 0
        cur.append(u)
        cur_len += len(u)
    if cur:
        groups.append(cur)

    # 3) 加句子级重叠：每组开头补上一组尾部累计 <= overlap 的句子
    out: list[str] = []
    for i, g in enumerate(groups):
        if i > 0 and overlap > 0:
            carry, clen = [], 0
            for u in reversed(groups[i - 1]):
                if clen + len(u) > overlap:
                    break
                carry.insert(0, u)
                clen += len(u)
            g = carry + g
        out.append(sep.join(g))

    # 短内容同样可能是有效知识（例如“测试”“同意”“OK”），不能因为低于
    # 最小块长度就把整篇文档过滤为空。对于长文末尾产生的短碎片，优先并入
    # 相邻 chunk；如果整篇只有一个短片段，则原样保留为一个 chunk。
    nonempty = [chunk.strip() for chunk in out if chunk.strip()]
    if not nonempty:
        return []

    min_len = 5 if cjk else 20
    chunks: list[str] = []
    leading_short = ""
    for chunk in nonempty:
        if len(chunk) < min_len:
            if chunks:
                chunks[-1] = f"{chunks[-1]}{sep}{chunk}".strip()
            else:
                leading_short = f"{leading_short}{sep}{chunk}".strip()
            continue

        if leading_short:
            chunk = f"{leading_short}{sep}{chunk}".strip()
            leading_short = ""
        chunks.append(chunk)

    if leading_short:
        chunks.append(leading_short)
    return chunks


def _ctx_prefix(title: str, heading_path: list) -> str:
    """把文档标题 + 当前标题层级拼成上下文前缀，如『出差标准 › 三、交通 › 3.1 飞机』。"""
    parts = [title] + [t for _, t in heading_path]
    return " › ".join(p for p in parts if p)


def _with_ctx(ctx: str, content: str) -> str:
    """给 chunk 内容加上标题路径上下文，提升检索命中率并让模型知道这段在讲什么。"""
    return f"【{ctx}】\n{content}" if ctx else content


def parse_markdown_content(content: str, source: str) -> list[dict]:
    """按标题层级切分，保留 Markdown 表格完整，并给每个 chunk 带上标题路径上下文。"""
    heading_re = re.compile(r'^(#{1,6})\s+(.*)$')
    table_row_re = re.compile(r'^\s*\|.*\|\s*$')

    chunks: list[dict] = []
    heading_path: list[tuple[int, str]] = []
    prose: list[str] = []
    table: list[str] = []
    discarded_placeholder_section = False

    def flush_prose():
        nonlocal discarded_placeholder_section
        if not prose:
            return
        raw_section = "\n".join(prose)
        prose.clear()
        if not raw_section.strip():
            return
        if _is_placeholder_only_section(raw_section):
            discarded_placeholder_section = True
            return
        section = raw_section.strip()
        ctx = _ctx_prefix(source, heading_path)
        for chunk in _split_text(section):
            chunks.append({"content": _with_ctx(ctx, chunk),
                           "metadata": {"source": source, "heading": ctx}})

    def flush_table():
        if not table:
            return
        md = "\n".join(table).strip()
        table.clear()
        if not md:
            return
        ctx = _ctx_prefix(source, heading_path)
        for piece in _split_markdown_table(md):
            chunks.append({"content": _with_ctx(ctx, piece),
                           "metadata": {"source": source, "heading": ctx, "type": "table"}})

    for line in content.split("\n"):
        h = heading_re.match(line)
        if h:
            flush_table()
            flush_prose()
            level = len(h.group(1))
            heading_path[:] = [x for x in heading_path if x[0] < level]
            text = h.group(2).strip()
            if text:
                heading_path.append((level, text))
        elif table_row_re.match(line):
            if prose:
                flush_prose()
            table.append(line)
        else:
            if table:
                flush_table()
            if line.strip() or prose:
                prose.append(line)

    flush_table()
    flush_prose()

    # 标题-only 文档仍使用原有兜底；但如果正文明确因“无/暂无”等占位符被
    # 丢弃，不能再把整篇原文切分回来，否则过滤会失效。
    if not chunks and content.strip() and not discarded_placeholder_section:
        for chunk in _split_text(content):
            chunks.append({"content": chunk, "metadata": {"source": source}})
    return chunks


def parse_file(filepath: str, source_name: str | None = None) -> list[dict]:
    ext = Path(filepath).suffix.lower()
    # source 用于 chunk 元数据与标题上下文前缀，优先用真实文件名（上传时保存路径是 uuid）
    source = source_name or os.path.basename(filepath)
    if ext == ".pdf":
        return _parse_pdf(filepath, source)
    elif ext == ".docx":
        return _parse_docx(filepath, source)
    elif ext == ".pptx":
        return _parse_pptx(filepath, source)
    elif ext in (".xlsx", ".xls"):
        return _parse_excel(filepath, source)
    else:
        return _parse_text(filepath, source)


def _parse_pdf(filepath: str, source: str) -> list[dict]:
    try:
        from pypdf import PdfReader
    except ImportError:
        from PyPDF2 import PdfReader  # 兼容旧依赖

    reader = PdfReader(filepath)
    title = os.path.splitext(source)[0]
    # 逐页处理：每页单独分块并带上「标题 › 第N页」上下文前缀，提升检索命中与可溯源性
    chunks: list[dict] = []
    for i, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        ctx = f"{title} › 第{i}页"
        for chunk in _split_text(text):
            chunks.append({"content": _with_ctx(ctx, chunk),
                           "metadata": {"source": source, "page": i, "heading": ctx}})
    return chunks


def _iter_docx_blocks(doc):
    """按文档原始顺序产出段落和表格（python-docx 的 doc.paragraphs 不含表格）。"""
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield "paragraph", Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield "table", Table(child, doc)


def _docx_cell_text(cell) -> str:
    text = "\n".join(p.text for p in cell.paragraphs).strip()
    # Markdown 表格单元格内不能有裸竖线和换行
    return text.replace("|", "\\|").replace("\n", "<br>")


def _docx_table_to_md(table) -> str:
    rows: list[list[str]] = []
    for row in table.rows:
        cells, prev_tc = [], None
        for cell in row.cells:
            # 水平合并会让 row.cells 重复返回同一个底层单元格，跳过重复
            if cell._tc is prev_tc:
                continue
            prev_tc = cell._tc
            cells.append(_docx_cell_text(cell))
        if any(c.strip() for c in cells):
            rows.append(cells)
    if not rows:
        return ""
    return _rows_to_md_table(rows[0], rows[1:])


def _split_markdown_table(md: str, max_chars: int = 1800) -> list[str]:
    """整张表尽量保留为一个 chunk；过大时按行拆分，每段重复表头。"""
    if len(md) <= max_chars:
        return [md]
    lines = md.split("\n")
    if len(lines) < 3:
        return [md]
    header, body = lines[:2], lines[2:]
    base = sum(len(h) for h in header)
    out, cur, cur_len = [], [], base
    for row in body:
        if cur and cur_len + len(row) > max_chars:
            out.append("\n".join(header + cur))
            cur, cur_len = [], base
        cur.append(row)
        cur_len += len(row)
    if cur:
        out.append("\n".join(header + cur))
    return out


def _docx_heading_level(paragraph) -> int | None:
    style = (paragraph.style.name if paragraph.style else "") or ""
    if style.lower().startswith("heading"):
        digits = "".join(ch for ch in style if ch.isdigit())
        return min(int(digits), 6) if digits else 1
    return None


def _parse_docx(filepath: str, source: str) -> list[dict]:
    from docx import Document
    doc = Document(filepath)
    title = os.path.splitext(source)[0]

    chunks: list[dict] = []
    heading_path: list[tuple[int, str]] = []
    text_buffer: list[str] = []

    def flush_text():
        if not text_buffer:
            return
        section = "\n".join(text_buffer)
        text_buffer.clear()
        ctx = _ctx_prefix(title, heading_path)
        for chunk in _split_text(section):
            chunks.append({"content": _with_ctx(ctx, chunk),
                           "metadata": {"source": source, "heading": ctx}})

    for kind, block in _iter_docx_blocks(doc):
        if kind == "paragraph":
            level = _docx_heading_level(block)
            if level is not None:
                # 新标题：先把上一节正文落盘（归到旧标题路径下），再更新路径
                flush_text()
                heading_path[:] = [h for h in heading_path if h[0] < level]
                if block.text.strip():
                    heading_path.append((level, block.text.strip()))
            elif block.text.strip():
                text_buffer.append(block.text)
        else:  # table：先把前面的正文落盘，再把表格作为完整 chunk 保留（带标题上下文）
            flush_text()
            md = _docx_table_to_md(block)
            if md:
                ctx = _ctx_prefix(title, heading_path)
                for piece in _split_markdown_table(md):
                    chunks.append({"content": _with_ctx(ctx, piece),
                                   "metadata": {"source": source, "heading": ctx, "type": "table"}})
    flush_text()
    return chunks


def docx_to_markdown(filepath: str) -> str:
    """把 docx 转成 Markdown（标题 + 表格）供预览/编辑，保留原文结构。"""
    from docx import Document
    doc = Document(filepath)
    parts: list[str] = []
    for kind, block in _iter_docx_blocks(doc):
        if kind == "paragraph":
            text = block.text.strip()
            if not text:
                continue
            style = (block.style.name if block.style else "") or ""
            if style.lower().startswith("heading"):
                digits = "".join(ch for ch in style if ch.isdigit())
                level = min(int(digits), 6) if digits else 1
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)
        else:
            md = _docx_table_to_md(block)
            if md:
                parts.append(md)
    return "\n\n".join(parts)


def _parse_pptx(filepath: str, source: str) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(filepath)
    title = os.path.splitext(source)[0]
    # 每张幻灯片作为结构单元单独分块，并带上「标题 › 第N页」上下文前缀
    chunks: list[dict] = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if not texts:
            continue
        ctx = f"{title} › 第{i}页"
        section = "\n".join(texts)
        for chunk in _split_text(section):
            chunks.append({"content": _with_ctx(ctx, chunk),
                           "metadata": {"source": source, "slide": i, "heading": ctx}})
    return chunks


def _rows_to_md_table(headers: list, rows: list) -> str:
    if not headers:
        return ""
    cols = len(headers)
    lines = ["| " + " | ".join(str(h) for h in headers) + " |"]
    lines.append("| " + " | ".join(["---"] * cols) + " |")
    for row in rows:
        padded = list(row) + [""] * max(0, cols - len(row))
        lines.append("| " + " | ".join(str(c) for c in padded[:cols]) + " |")
    return "\n".join(lines)


def excel_to_markdown(filepath: str) -> str:
    """Convert Excel file to Markdown tables for preview/editing."""
    ext = Path(filepath).suffix.lower()
    parts = []
    if ext == ".xls":
        import xlrd
        wb = xlrd.open_workbook(filepath)
        for i in range(wb.nsheets):
            ws = wb.sheet_by_index(i)
            if ws.nrows == 0:
                continue
            headers = [str(c).strip() for c in ws.row_values(0)]
            rows = [[str(c).strip() for c in ws.row_values(r)] for r in range(1, ws.nrows)]
            parts.append(f"## {ws.name}\n\n{_rows_to_md_table(headers, rows)}")
    else:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_rows = list(ws.iter_rows(values_only=True))
            if not all_rows:
                continue
            headers = [str(c).strip() if c is not None else "" for c in all_rows[0]]
            rows = [[str(c).strip() if c is not None else "" for c in row] for row in all_rows[1:]]
            parts.append(f"## {sheet_name}\n\n{_rows_to_md_table(headers, rows)}")
        wb.close()
    return "\n\n".join(parts)


def _parse_excel(filepath: str, source: str) -> list[dict]:
    ext = Path(filepath).suffix.lower()
    return _excel_to_chunks(filepath, source, ext)


def _excel_sheet_to_text(headers: list, rows: list[list]) -> str:
    lines = []
    for row in rows:
        pairs = [f"{h}: {v}" for h, v in zip(headers, row) if str(v).strip()]
        if pairs:
            lines.append("，".join(pairs))
    return "\n".join(lines)


def _excel_to_chunks(filepath: str, source: str, ext: str) -> list[dict]:
    chunks: list[dict] = []

    if ext == ".xls":
        import xlrd
        wb = xlrd.open_workbook(filepath)
        for i in range(wb.nsheets):
            ws = wb.sheet_by_index(i)
            if ws.nrows < 2:
                continue
            headers = [str(c).strip() for c in ws.row_values(0)]
            rows = [[str(c).strip() for c in ws.row_values(r)] for r in range(1, ws.nrows)]
            text = _excel_sheet_to_text(headers, rows)
            if text:
                section = f"[{ws.name}]\n{text}"
                for chunk in _split_text(section):
                    chunks.append({"content": chunk, "metadata": {"source": source, "sheet": ws.name}})
    else:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_rows = list(ws.iter_rows(values_only=True))
            if len(all_rows) < 2:
                continue
            headers = [str(c).strip() if c is not None else "" for c in all_rows[0]]
            rows = [[str(c).strip() if c is not None else "" for c in row] for row in all_rows[1:]]
            text = _excel_sheet_to_text(headers, rows)
            if text:
                section = f"[{sheet_name}]\n{text}"
                for chunk in _split_text(section):
                    chunks.append({"content": chunk, "metadata": {"source": source, "sheet": sheet_name}})
        wb.close()

    return chunks


def _parse_text(filepath: str, source: str) -> list[dict]:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        full_text = f.read()
    # 走 Markdown 结构化解析：有标题按标题层级切、保留表格，并给每块加标题路径上下文；
    # 纯文本无标题时也会带上文件名作为上下文前缀。
    return parse_markdown_content(full_text, source)
